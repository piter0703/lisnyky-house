#!/usr/bin/env python3
"""
Art History Test Answers Presentation Generator
Creates a professional dark-themed PPTX with answers to all 20 questions.
"""

import os
import io
import time
import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from PIL import Image as PILImage

# ── COLORS ────────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x0D, 0x0D)   # near-black
GOLD     = RGBColor(0xE8, 0xA0, 0x20)   # golden orange (like the PDF)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xCC, 0xCC, 0xCC)
DGRAY    = RGBColor(0x33, 0x33, 0x33)
GREEN    = RGBColor(0x4C, 0xAF, 0x50)
ACCBLUE  = RGBColor(0x42, 0x9E, 0xCB)

# ── SLIDE DIMENSIONS (16:9) ────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

IMG_DIR = "/tmp/art_pptx_imgs"
os.makedirs(IMG_DIR, exist_ok=True)

HEADERS = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36"}

def download_img(url: str, name: str) -> str | None:
    path = os.path.join(IMG_DIR, name)
    if os.path.exists(path):
        return path
    try:
        r = requests.get(url, headers=HEADERS, timeout=15)
        if r.status_code == 200 and len(r.content) > 1000:
            # convert to JPEG for pptx compatibility
            img = PILImage.open(io.BytesIO(r.content)).convert("RGB")
            img.save(path, "JPEG", quality=88)
            return path
    except Exception as e:
        print(f"  ⚠ Could not download {name}: {e}")
    return None

# ── IMAGE URLS (Wikimedia Commons – public domain) ────────────────────────────
IMAGES = {
    # Q1 – Paleolithic Venuses
    "venus_willendorf":  "https://upload.wikimedia.org/wikipedia/commons/thumb/5/50/Venus_von_Willendorf_01.jpg/600px-Venus_von_Willendorf_01.jpg",
    "venus_laussel":     "https://upload.wikimedia.org/wikipedia/commons/thumb/c/ca/Venus_of_Laussel.jpg/600px-Venus_of_Laussel.jpg",
    "venus_lespugue":    "https://upload.wikimedia.org/wikipedia/commons/thumb/0/08/VenusLespugue.jpg/600px-VenusLespugue.jpg",

    # Q2 – Egyptian burial
    "mummy":             "https://upload.wikimedia.org/wikipedia/commons/thumb/4/4e/Mummy_of_Ramesses_II.jpg/600px-Mummy_of_Ramesses_II.jpg",
    "canopic_jars":      "https://upload.wikimedia.org/wikipedia/commons/thumb/e/e7/Canopic_jar_lids.jpg/600px-Canopic_jar_lids.jpg",
    "sarcophagus":       "https://upload.wikimedia.org/wikipedia/commons/thumb/e/e8/Tut_gold_coffin.jpg/600px-Tut_gold_coffin.jpg",
    "ankh":              "https://upload.wikimedia.org/wikipedia/commons/thumb/a/ae/Ankh_em_tut.jpg/600px-Ankh_em_tut.jpg",
    "death_mask":        "https://upload.wikimedia.org/wikipedia/commons/thumb/2/27/Tutankhamun_mask.jpg/600px-Tutankhamun_mask.jpg",
    "ushabti":           "https://upload.wikimedia.org/wikipedia/commons/thumb/9/9c/Ushabti_Louvre.jpg/600px-Ushabti_Louvre.jpg",
    "ba_soul":           "https://upload.wikimedia.org/wikipedia/commons/thumb/8/8d/Ba_bird.jpg/600px-Ba_bird.jpg",

    # Q3 – Egyptian statues
    "hesyra":            "https://upload.wikimedia.org/wikipedia/commons/thumb/4/4b/Hesy-Ra_CG_1426.jpg/600px-Hesy-Ra_CG_1426.jpg",
    "kaaper":            "https://upload.wikimedia.org/wikipedia/commons/thumb/2/28/Egyptian_-_Standing_Male_Figure_%28the_Sheikh_el-Beled%29_-_Walters_22235_-_Three_Quarter_Right.jpg/600px-Egyptian_-_Standing_Male_Figure_%28the_Sheikh_el-Beled%29_-_Walters_22235_-_Three_Quarter_Right.jpg",
    "kai_scribe":        "https://upload.wikimedia.org/wikipedia/commons/thumb/b/b5/Squatting_scribe_Louvre_051906.jpg/600px-Squatting_scribe_Louvre_051906.jpg",
    "seneb":             "https://upload.wikimedia.org/wikipedia/commons/thumb/4/4f/Seneb_and_family.jpg/600px-Seneb_and_family.jpg",
    "rahotep_nofret":    "https://upload.wikimedia.org/wikipedia/commons/thumb/0/0c/Egyptian_Museum%2C_Cairo_-_Rahotep_and_Nofret.jpg/600px-Egyptian_Museum%2C_Cairo_-_Rahotep_and_Nofret.jpg",

    # Q4 – Aegean art
    "la_parisienne":     "https://upload.wikimedia.org/wikipedia/commons/thumb/5/57/Minoan_woman_or_goddess.jpg/600px-Minoan_woman_or_goddess.jpg",
    "dolphin_fresco":    "https://upload.wikimedia.org/wikipedia/commons/thumb/a/ae/Akrotiri_dolphins.jpg/600px-Akrotiri_dolphins.jpg",
    "prince_lilies":     "https://upload.wikimedia.org/wikipedia/commons/thumb/1/12/Knossos_Lilienprinz.jpg/600px-Knossos_Lilienprinz.jpg",
    "blue_ladies":       "https://upload.wikimedia.org/wikipedia/commons/thumb/d/d2/Blue_ladies.jpg/600px-Blue_ladies.jpg",

    # Q5 – Greek vessels
    "amphora":           "https://upload.wikimedia.org/wikipedia/commons/thumb/4/47/Amphora_Exekias_BM_B210.jpg/600px-Amphora_Exekias_BM_B210.jpg",
    "hydria":            "https://upload.wikimedia.org/wikipedia/commons/thumb/1/1a/Hydria_Meidias_Painter_BM_E224.jpg/600px-Hydria_Meidias_Painter_BM_E224.jpg",
    "krater":            "https://upload.wikimedia.org/wikipedia/commons/thumb/b/bb/Calyx_krater_Euphronios_Louvre.jpg/600px-Calyx_krater_Euphronios_Louvre.jpg",
    "kantharos":         "https://upload.wikimedia.org/wikipedia/commons/thumb/e/e4/Kantharos_Musee_Rodin.jpg/600px-Kantharos_Musee_Rodin.jpg",
    "kylix":             "https://upload.wikimedia.org/wikipedia/commons/thumb/1/1b/Kylix_warrior_Louvre_G128.jpg/600px-Kylix_warrior_Louvre_G128.jpg",

    # Q6 – Entrance structures
    "lion_gate":         "https://upload.wikimedia.org/wikipedia/commons/thumb/b/b6/Lions_gate_Mycenae_a.jpg/600px-Lions_gate_Mycenae_a.jpg",
    "propylaea":         "https://upload.wikimedia.org/wikipedia/commons/thumb/2/28/Propylaea_of_Parthenon.jpg/600px-Propylaea_of_Parthenon.jpg",
    "triumph_arch":      "https://upload.wikimedia.org/wikipedia/commons/thumb/4/4e/Arch_Septimius_Severus.jpg/600px-Arch_Septimius_Severus.jpg",

    # Q7 – Ancient Greek sculpture
    "apollo_belvedere":  "https://upload.wikimedia.org/wikipedia/commons/thumb/5/50/Apollo_of_the_Belvedere.jpg/600px-Apollo_of_the_Belvedere.jpg",
    "venus_milo":        "https://upload.wikimedia.org/wikipedia/commons/thumb/2/26/Venus_de_Milo_Louvre_Ma399_n4.jpg/600px-Venus_de_Milo_Louvre_Ma399_n4.jpg",
    "nike_samothrace":   "https://upload.wikimedia.org/wikipedia/commons/thumb/b/b5/Nike_of_Samothrake_Louvre_Ma2369_n4.jpg/600px-Nike_of_Samothrake_Louvre_Ma2369_n4.jpg",
    "hermes_dionysus":   "https://upload.wikimedia.org/wikipedia/commons/thumb/7/7a/Hermes_Praxiteles_Archaeological_Museum_Olympia.jpg/600px-Hermes_Praxiteles_Archaeological_Museum_Olympia.jpg",
    "laocoon":           "https://upload.wikimedia.org/wikipedia/commons/thumb/b/b6/Laocoon_and_His_Sons.jpg/600px-Laocoon_and_His_Sons.jpg",
    "maenad":            "https://upload.wikimedia.org/wikipedia/commons/thumb/e/e8/Skopas_maenad.jpg/600px-Skopas_maenad.jpg",

    # Q8 – Roman portraits
    "augustus_prima":    "https://upload.wikimedia.org/wikipedia/commons/thumb/e/eb/Statue-Augustus.jpg/600px-Statue-Augustus.jpg",
    "roman_ancestor":    "https://upload.wikimedia.org/wikipedia/commons/thumb/5/52/Barberini_Togatus.jpg/600px-Barberini_Togatus.jpg",
    "marcus_aurelius":   "https://upload.wikimedia.org/wikipedia/commons/thumb/3/35/Equestrian_statue_of_Marcus_Aurelius_Musei_Capitolini_MC4.jpg/600px-Equestrian_statue_of_Marcus_Aurelius_Musei_Capitolini_MC4.jpg",
    "antinous":          "https://upload.wikimedia.org/wikipedia/commons/thumb/6/6e/Antinoos_Mondragone_Louvre_Ma1205.jpg/600px-Antinoos_Mondragone_Louvre_Ma1205.jpg",
    "roman_lady":        "https://upload.wikimedia.org/wikipedia/commons/thumb/9/98/Head_of_a_Roman_Lady_from_the_Flavian_period.jpg/600px-Head_of_a_Roman_Lady_from_the_Flavian_period.jpg",

    # Q9 – Medieval architecture
    "hagia_sophia":      "https://upload.wikimedia.org/wikipedia/commons/thumb/2/22/Hagia_Sophia_Istanbul_2014_2.jpg/800px-Hagia_Sophia_Istanbul_2014_2.jpg",
    "worms_cathedral":   "https://upload.wikimedia.org/wikipedia/commons/thumb/7/78/Worms_Dom_BW_1.jpg/600px-Worms_Dom_BW_1.jpg",
    "chartres":          "https://upload.wikimedia.org/wikipedia/commons/thumb/7/70/Chartres_Cathedral_North_Portal.jpg/600px-Chartres_Cathedral_North_Portal.jpg",
    "notre_dame":        "https://upload.wikimedia.org/wikipedia/commons/thumb/b/b6/Notre-Dame_de_Paris_2013-07-24.jpg/600px-Notre-Dame_de_Paris_2013-07-24.jpg",

    # Q10 – Gothic facade elements  (same as Amiens)
    "amiens_full":       "https://upload.wikimedia.org/wikipedia/commons/thumb/8/8a/Amiens_Cathédrale_Notre-Dame_Fassade_2.jpg/600px-Amiens_Cathédrale_Notre-Dame_Fassade_2.jpg",

    # Q13 – Byzantine icons
    "byzantine_mosaic":  "https://upload.wikimedia.org/wikipedia/commons/thumb/a/a3/Ravenna_mosaic_Justinian.jpg/800px-Ravenna_mosaic_Justinian.jpg",

    # Q15 – Virgin Mary types
    "orans_madonna":     "https://upload.wikimedia.org/wikipedia/commons/thumb/1/1f/Hagia_Sophia_-_Apse_mosaic_2.jpg/600px-Hagia_Sophia_-_Apse_mosaic_2.jpg",
    "hodegetria":        "https://upload.wikimedia.org/wikipedia/commons/thumb/d/de/Our_Lady_of_Kazan.jpg/600px-Our_Lady_of_Kazan.jpg",
    "eleusa":            "https://upload.wikimedia.org/wikipedia/commons/thumb/e/e6/Vladimir_mother_of_god_from_tretyakov.jpg/600px-Vladimir_mother_of_god_from_tretyakov.jpg",

    # Q17 – David in Florence
    "david_michelangelo": "https://upload.wikimedia.org/wikipedia/commons/thumb/8/80/Michelangelo%27s_David_-_right_view_2.jpg/600px-Michelangelo%27s_David_-_right_view_2.jpg",
    "david_donatello":   "https://upload.wikimedia.org/wikipedia/commons/thumb/7/76/Donatello%2C_David_with_the_Head_of_Goliath%2C_1440s%2C_Bargello%2C_Florence.jpg/600px-Donatello%2C_David_with_the_Head_of_Goliath%2C_1440s%2C_Bargello%2C_Florence.jpg",
    "david_verrocchio":  "https://upload.wikimedia.org/wikipedia/commons/thumb/4/4f/Verrocchio_David_Bargello.jpg/600px-Verrocchio_David_Bargello.jpg",

    # Q18 – Renaissance techniques
    "sistine_fresco":    "https://upload.wikimedia.org/wikipedia/commons/thumb/5/5b/Michelangelo_-_Creation_of_Adam_%28cropped%29.jpg/800px-Michelangelo_-_Creation_of_Adam_%28cropped%29.jpg",
    "mona_lisa":         "https://upload.wikimedia.org/wikipedia/commons/thumb/e/ec/Mona_Lisa%2C_by_Leonardo_da_Vinci%2C_from_C2RMF_retouched.jpg/600px-Mona_Lisa%2C_by_Leonardo_da_Vinci%2C_from_C2RMF_retouched.jpg",
    "botticelli_simonetta": "https://upload.wikimedia.org/wikipedia/commons/thumb/b/b1/Simonetta_Vespucci_by_Sandro_Botticelli.jpg/600px-Simonetta_Vespucci_by_Sandro_Botticelli.jpg",

    # Q19 – Mythology paintings
    "botticelli_spring": "https://upload.wikimedia.org/wikipedia/commons/thumb/3/3c/Botticelli-primavera.jpg/1200px-Botticelli-primavera.jpg",
    "botticelli_venus":  "https://upload.wikimedia.org/wikipedia/commons/thumb/2/26/Sandro_Botticelli_-_La_nascita_di_Venere_-_Google_Art_Project_-_edited.jpg/1200px-Sandro_Botticelli_-_La_nascita_di_Venere_-_Google_Art_Project_-_edited.jpg",
    "titian_europa":     "https://upload.wikimedia.org/wikipedia/commons/thumb/4/44/Titian_-_The_Rape_of_Europa_-_Google_Art_Project.jpg/800px-Titian_-_The_Rape_of_Europa_-_Google_Art_Project.jpg",
    "botticelli_venus_mars": "https://upload.wikimedia.org/wikipedia/commons/thumb/1/15/Botticelli-Venus-and-Mars.jpg/800px-Botticelli-Venus-and-Mars.jpg",

    # Q20 – Renaissance works
    "fra_angelico_annunciation": "https://upload.wikimedia.org/wikipedia/commons/thumb/b/b4/Annunciation_%28Fra_Angelico%2C_1432%29.jpg/800px-Annunciation_%28Fra_Angelico%2C_1432%29.jpg",
    "piero_diptych":     "https://upload.wikimedia.org/wikipedia/commons/thumb/2/2f/Piero_della_Francesca_-_Diptych_of_the_Dukes_of_Urbino_%28obverses%29_-_WGA17609.jpg/800px-Piero_della_Francesca_-_Diptych_of_the_Dukes_of_Urbino_%28obverses%29_-_WGA17609.jpg",
    "masaccio_tribute":  "https://upload.wikimedia.org/wikipedia/commons/thumb/3/3a/Masaccio_-_Tribute_Money_-_WGA14254.jpg/800px-Masaccio_-_Tribute_Money_-_WGA14254.jpg",
    "mantegna_dead_christ": "https://upload.wikimedia.org/wikipedia/commons/thumb/3/32/Mantegna_dead_christ.jpg/800px-Mantegna_dead_christ.jpg",
    "titian_venus_urbino": "https://upload.wikimedia.org/wikipedia/commons/thumb/2/2c/Tiziano_-_Venere_di_Urbino_-_Google_Art_Project.jpg/800px-Tiziano_-_Venere_di_Urbino_-_Google_Art_Project.jpg",
    "bellini_madonna":   "https://upload.wikimedia.org/wikipedia/commons/thumb/9/9a/Giovanni_Bellini%2C_Madonna_of_the_Meadow_-_National_Gallery.jpg/600px-Giovanni_Bellini%2C_Madonna_of_the_Meadow_-_National_Gallery.jpg",
}

# ── PPT HELPERS ───────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_image(slide, path, x, y, w, h):
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, x, y, w, h)
            return True
        except Exception as e:
            print(f"  ⚠ Could not add image {path}: {e}")
    return False

def title_slide(prs, title, subtitle=""):
    slide = blank_slide(prs)
    fill_bg(slide)
    # decorative bar top
    add_rect(slide, 0, 0, W, Inches(0.08), GOLD)
    # decorative bar bottom
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), GOLD)
    # title
    add_text(slide, title, Inches(1), Inches(2.5), Inches(11.33), Inches(1.8),
             font_size=42, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, Inches(1), Inches(4.5), Inches(11.33), Inches(1),
                 font_size=22, color=LGRAY, align=PP_ALIGN.CENTER)
    return slide

def section_slide(prs, section_name):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.08), GOLD)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), GOLD)
    add_rect(slide, Inches(1), Inches(2.8), Inches(11.33), Inches(2), DGRAY)
    add_text(slide, section_name, Inches(1.2), Inches(3.0), Inches(11), Inches(1.6),
             font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    return slide

def answer_slide_1col(prs, q_num, question, answers, img_keys=None, notes=""):
    """Slide with text on left, images on right."""
    slide = blank_slide(prs)
    fill_bg(slide)
    # top bar
    add_rect(slide, 0, 0, W, Inches(0.06), GOLD)
    # question number badge
    add_rect(slide, Inches(0.15), Inches(0.15), Inches(0.65), Inches(0.55), GOLD)
    add_text(slide, str(q_num), Inches(0.15), Inches(0.12), Inches(0.65), Inches(0.55),
             font_size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)
    # question title
    add_text(slide, question, Inches(0.9), Inches(0.12), Inches(10), Inches(0.6),
             font_size=15, bold=True, color=GOLD)

    # separator line
    add_rect(slide, Inches(0.15), Inches(0.8), Inches(12.8), Inches(0.03), DGRAY)

    # answers block (left side)
    text_w = Inches(6.0) if img_keys else Inches(12.5)
    y_start = Inches(0.95)
    full_text = "\n".join(answers)
    add_text(slide, full_text, Inches(0.2), y_start, text_w, Inches(5.8),
             font_size=13, color=WHITE, wrap=True)

    # images (right side, up to 4)
    if img_keys:
        imgs = [download_img(IMAGES[k], k + ".jpg") for k in img_keys if k in IMAGES]
        imgs = [i for i in imgs if i]
        n = len(imgs)
        if n == 1:
            add_image(slide, imgs[0], Inches(6.4), Inches(0.9), Inches(6.5), Inches(5.5))
        elif n == 2:
            add_image(slide, imgs[0], Inches(6.4), Inches(0.9), Inches(3.1), Inches(5.4))
            add_image(slide, imgs[1], Inches(9.8), Inches(0.9), Inches(3.1), Inches(5.4))
        elif n == 3:
            add_image(slide, imgs[0], Inches(6.4), Inches(0.9), Inches(2.0), Inches(5.4))
            add_image(slide, imgs[1], Inches(8.6), Inches(0.9), Inches(2.0), Inches(5.4))
            add_image(slide, imgs[2], Inches(10.8), Inches(0.9), Inches(2.1), Inches(5.4))
        elif n >= 4:
            add_image(slide, imgs[0], Inches(6.4), Inches(0.9), Inches(3.1), Inches(2.6))
            add_image(slide, imgs[1], Inches(9.8), Inches(0.9), Inches(3.1), Inches(2.6))
            add_image(slide, imgs[2], Inches(6.4), Inches(3.7), Inches(3.1), Inches(2.6))
            add_image(slide, imgs[3], Inches(9.8), Inches(3.7), Inches(3.1), Inches(2.6))

    # optional notes at bottom
    if notes:
        add_rect(slide, 0, H - Inches(0.55), W, Inches(0.55), DGRAY)
        add_text(slide, notes, Inches(0.2), H - Inches(0.52), Inches(12.8), Inches(0.5),
                 font_size=10, color=LGRAY, italic=True)
    return slide


# ── DOWNLOAD ALL IMAGES FIRST ─────────────────────────────────────────────────
def download_all():
    print("📥 Downloading images…")
    ok = 0
    for name, url in IMAGES.items():
        path = download_img(url, name + ".jpg")
        if path:
            ok += 1
        time.sleep(0.3)
    print(f"   {ok}/{len(IMAGES)} images ready\n")

# ── BUILD PRESENTATION ────────────────────────────────────────────────────────
def build():
    download_all()

    prs = new_prs()

    # ── TITLE ──────────────────────────────────────────────────────────────────
    title_slide(prs,
                "МИСТЕЦТВО ДАВНЬОГО СВІТУ\nТА СЕРЕДНІХ ВІКІВ",
                "Відповіді на тестові завдання · 2026")

    # ══════════════════════════════════════════════════════════════════════════
    section_slide(prs, "МИСТЕЦТВО ДАВНЬОГО СВІТУ")

    # ── Q1 ─────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 1,
        "Оберіть палеолітичні венери",
        [
            "✅  ПРАВИЛЬНІ ВІДПОВІДІ: 3, 5 (і, можливо, 8)",
            "",
            "Палеолітична венера — жіноча фігурка, ліплена 25 000–15 000 р. до н.е.",
            "Характерні ознаки:",
            "  • Гіпертрофовані груди, стегна, живіт — символ родючості",
            "  • Відсутність обличчя або рис обличчя",
            "  • Маленькі ноги та руки",
            "  • Матеріал: вапняк, мамонтова кістка",
            "",
            "№ 5 — Венера Віллендорфська (~24 000 р. до н.е., Австрія)",
            "№ 3 — Венера Лоссельська (барельєф, ~22 000 р. до н.е., Франція)",
            "№ 8 — Спляча жінка з Гозо (неоліт, Мальта) — умовно подібна",
            "",
            "Решта зображень (#1, 2, 4, 6, 7) — класична антична скульптура.",
        ],
        ["venus_willendorf", "venus_laussel"],
        "Венера Віллендорфська | Венера Лоссельська"
    )

    # ── Q2 ─────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 2,
        "Поховальний культ давніх єгиптян — відповідність",
        [
            "А) мумія              → фото 5",
            "Б) канопи             → фото 4  (глеки з головами богів)",
            "В) саркофаг           → фото 1  (золота труна Тутанхамона)",
            "Г) амулет анх         → фото 8  (хрест із петлею = символ життя)",
            "Д) поховальна маска   → фото 2  (золота маска Тутанхамона)",
            "Е) ушебті             → фото 3  (маленькі фігурки-слуги)",
            "Є) душа «ба»          → фото 6  (птах із людською головою)",
            "Ж) дух «ка»           → фото 9  (фігура з піднятими руками)",
        ],
        ["sarcophagus", "death_mask", "canopic_jars", "ushabti"],
        "саркофаг · маска · канопи · ушебті"
    )

    # ── Q3 ─────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 3,
        "Давньоєгипетські статуї та рельєфи",
        [
            "А) Архітектор Хесіра  → фото 2",
            "   Дерев'яний рельєф із мастаби, бл. 2650 р. до н.е.",
            "",
            "Б) Жрець Каапер       → фото 1",
            "   «Сільський старійшина» — реалістична дерев'яна статуя, ≈2500 р. до н.е.",
            "",
            "В) Писар Каї          → фото 6",
            "   Сидяча постава «лотос», схрещені ноги, сувій на колінах, Луврський музей",
            "",
            "Г) Карлик Сенеб       → фото 5",
            "   Сидить з дружиною, діти розміщені там, де в повних людей мали б бути ноги",
            "",
            "Д) Рахотеп і Нофрет  → фото 4 (пара кольорових статуй, ок. 2600 р. до н.е.)",
        ],
        ["hesyra", "kaaper", "kai_scribe", "seneb", "rahotep_nofret"],
        "Єгипетська скульптура досягала вражаючого психологічного реалізму"
    )

    # ── Q4 ─────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 4,
        "Егейське (крито-мікенське) мистецтво — розписи",
        [
            "✅  ПРАВИЛЬНІ ВІДПОВІДІ: 1, 3, 4, 5, 7",
            "",
            "№ 1 — «Ла Парізьєнка» — фреска з Кноссу, ≈1500 р. до н.е.",
            "№ 3 — Фреска з дельфінами, палац Кносс",
            "№ 4 — «Блакитні дами», жрецькі процесії",
            "№ 5 — «Принц лілій» — фреска з Кноссу, хлопець у квітах",
            "№ 7 — Фреска з носіями ваз (Кіпр/Крит)",
            "",
            "❌  НЕ КРИТСЬКІ:",
            "№ 2 — єгипетська садова фреска",
            "№ 6 — еллінська/римська",
            "№ 8 — єгипетські музикантки",
            "",
            "Ознаки крито-мікенського мистецтва: жива лінія, рух, синьо-червона палітра",
        ],
        ["la_parisienne", "dolphin_fresco", "prince_lilies"],
        "Мінойське мистецтво — найбільш «живе» з усіх стародавніх цивілізацій"
    )

    # ── Q5 ─────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 5,
        "Давньогрецькі сосуди — типи і призначення",
        [
            "А) АМФОРА    → фото 1, 10",
            "   Призначення: зберігання вина, оливкової олії, зерна",
            "   Ознаки: 2 вертикальні ручки, конічне дно або підніжка",
            "",
            "Б) ГІДРІЯ    → фото 4, 9",
            "   Призначення: носіння і зберігання води",
            "   Ознаки: 3 ручки (дві горизонт. + одна вертикальна)",
            "",
            "В) КРАТЕР    → фото 5",
            "   Призначення: змішування вина з водою під час банкету",
            "   Ознаки: широкий посуд із великим горлом",
            "",
            "Г) КАНФАР    → фото 2",
            "   Призначення: питна кубок для вина, атрибут Діоніса",
            "   Ознаки: високі вигнуті ручки, що піднімаються вище вінця",
            "",
            "Д) КІЛІК     → фото 6, 8",
            "   Призначення: питна чаша для вина (симпозіум)",
            "   Ознаки: широка плоска чаша на ніжці, дві горизонтальні ручки",
        ],
        ["amphora", "kantharos", "krater", "kylix"],
        "Фото 3, 7 — ойнохоя (глечик для вина); фото не відповідає жодному з А–Д"
    )

    # ── Q6 ─────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 6,
        "Архітектурні входи — відповідність",
        [
            "А) «Левові ворота» → фото 1",
            "   Мікени, Греція, ≈1250 р. до н.е.",
            "   Провідні: у цитадель (акрополь) міста Мікени",
            "",
            "Б) Пропілеї → фото 3",
            "   Афінський акрополь, 437–432 рр. до н.е., арх. Мнесікл",
            "   Провідні: на священну скелю Акрополя",
            "",
            "В) Тріумфальна арка → фото 5",
            "   Арка Септимія Севера, Рим, 203 р. н.е.",
            "   Провідні: до Римського Форуму",
            "",
            "НЕ ВХОДИ (додаткові):",
            "№ 2 — Абу-Симбел, Єгипет (фасад скельного храму)",
            "№ 4 — Ерехтейон, каріатиди (портик, не вхід)",
            "№ 6 — Луксорський пілон (вхід до єгипетського храму)",
            "№ 7 — Міст Гар (Pont du Gard), Франція — акведук",
        ],
        ["lion_gate", "propylaea", "triumph_arch"],
        "Пропілеї — парадний вхід; Левові ворота — найдавніший монументальний вхід у Європі"
    )

    # ── Q7 ─────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 7,
        "Античні скульптури — відповідність + період",
        [
            "А) «Танцююча менада» → фото 6",
            "   Автор: Скопас, ≈350 р. до н.е. — КЛАСИКА",
            "   Перша скульптура з відкинутою головою, пристрастю в русі",
            "",
            "Б) «Гермес з немовлям Діонісом» → фото 4",
            "   Автор: Пракситель, ≈330 р. до н.е. — КЛАСИКА",
            "   М'яка моделировка тіла, S-подібний вигин",
            "",
            "В) «Аполлон Бельведерський» → фото 1",
            "   Римська копія, оригінал бл. 330–325 р. до н.е. — КЛАСИКА",
            "",
            "Г) «Венера Мілоська» → фото 2",
            "   Агесандр (або невідомий), ≈130–100 р. до н.е. — ЕЛЛІНІЗМ",
            "   Гострий ритм складок, контраст оголеного тіла і задрапованого низу",
            "",
            "Д) «Лаокоон та його сини» → фото 5",
            "   Родоські скульптори, ≈25 р. до н.е. — ЕЛЛІНІЗМ",
            "   Вершина драматизму і психологічного напруження",
            "",
            "№ 3 — Ніка Самофракійська — ЕЛЛІНІЗМ (≈190 р. до н.е.)",
        ],
        ["apollo_belvedere", "venus_milo", "hermes_dionysus", "laocoon"],
        "Класика = гармонія та ідеал; Еллінізм = емоція, рух, драматизм"
    )

    # ── Q8 ─────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 8,
        "Давньоримські скульптурні портрети",
        [
            "А) Римлянин з бюстами предків → фото 3",
            "   «Барберіні тогатус» — I ст. до н.е.",
            "   Традиція ius imaginum: право носити воскові маски предків",
            "",
            "Б) Імператор Октавіан Август з Пріма Порта → фото 1",
            "   Мармур, ≈20 р. до н.е. Тип ідеалізованого героя + Купідон біля ніг",
            "",
            "В) Бюст Антиноя в образі Діоніса → фото 6",
            "   Улюбленець Адріана, обожнений після смерті у 130 р. н.е.",
            "",
            "Г) Кінний монумент Марка Аврелія → фото 4",
            "   Бронза, ≈165–170 р. н.е. Єдиний античний кінний монумент, що зберігся",
            "",
            "Д) Бюст римлянки часів Флавіїв → фото 7",
            "   II пол. I ст. н.е. Характерна: складна зачіска з локонами «пагода»",
        ],
        ["augustus_prima", "roman_ancestor", "marcus_aurelius", "antinous"],
        "Рим = портретний реалізм + ідеалізація влади; єдине поєднання в античному мистецтві"
    )

    # ══════════════════════════════════════════════════════════════════════════
    section_slide(prs, "МИСТЕЦТВО СЕРЕДНІХ ВІКІВ")

    # ── Q9 ─────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 9,
        "Середньовічні архітектурні стилі",
        [
            "А) ВІЗАНТІЙСЬКИЙ стиль → фото 6",
            "   Собор Святої Софії, Стамбул (Константинополь), 532–537 рр.",
            "   Ознаки: купол над квадратом, золоті мозаїки, світло через вікна барабана",
            "",
            "Б) РОМАНСЬКИЙ стиль → фото 7, 8",
            "   фото 7 — замок Шільон, Швейцарія",
            "   фото 8 — Вормський собор, Німеччина, XI–XII ст.",
            "   Ознаки: товсті стіни, напівциркульні арки, темні інтер'єри",
            "",
            "В) ГОТИЧНИЙ стиль → фото 1, 5",
            "   фото 1 — собор Нотр-Дам де Шартр, Франція",
            "   фото 5 — Нотр-Дам де Парі",
            "   Ознаки: стрільчасті арки, контрфорси, вітражі, вертикаль",
            "",
            "НЕ ВІДПОВІДІ:",
            "фото 2 — Андріївська церква, Київ (бароко, XVIII ст.)",
            "фото 3 — Каркасон (романська фортеця)",
            "фото 4 — модерн (сецесія), кін. XIX ст.",
        ],
        ["hagia_sophia", "worms_cathedral", "chartres", "notre_dame"],
        "Гасло готики: «Камінь, що тягнеться до неба»"
    )

    # ── Q10 ────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 10,
        "Архітектурні елементи готичного собору (Ам'єн)",
        [
            "А) БАШТИ → позначені «3» (ліворуч і праворуч угорі)",
            "   Дві масивні вежі фасаду, завершені шпилями — символ вертикалі",
            "",
            "Б) ПЕРСПЕКТИВНИЙ ПОРТАЛ → позначений «1» (внизу по центру)",
            "   Ступінчасте заглиблення з рядами статуй; передає глибину простору",
            "",
            "В) ГОТИЧНА ТРОЯНДА (РОЗА) → позначена «5» (великий круглий вітраж)",
            "   Кам'яна ажурна решітка із кольоровим склом, символ Богоматері",
            "",
            "Г) СТРІЛЬЧАСТІ АРКИ → позначені «4» (арки порталів)",
            "   Загострені догори, знімають бічний тиск стін — ключ готичної конструкції",
            "",
            "Д) ГАЛЕРЕЯ КОРОЛІВ → позначена «2»",
            "   Горизонтальний ряд статуй між порталами і трояндою",
            "   28 статуй юдейських царів (предків Христа)",
        ],
        ["amiens_full"],
        "Фасад собору в Ам'єні — «камінна енциклопедія» середньовіччя"
    )

    # ── Q11 ────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 11,
        "Капітелі середньовічних колон",
        [
            "А) КУБІЧНА КАПІТЕЛЬ → фото 6 — РОМАНСЬКИЙ",
            "   Проста кубічна форма, нижня частина заокруглена",
            "   Прямий наступник давньоримської дорики",
            "",
            "Б) З МОТИВОМ ПЛЕТІНКИ → фото 3 — РОМАНСЬКИЙ",
            "   Геометричний орнамент, плетений візерунок",
            "   Характерний для ранньохристиянського і романського мистецтва",
            "",
            "В) «АДАМ І ЄВА» → фото 2 — РОМАНСЬКИЙ",
            "   Наративний сюжет, схематичні фігури, відсутність перспективи",
            "",
            "Г) З ХИМЕРНИМИ ОБРАЗАМИ → фото 4 — РОМАНСЬКИЙ",
            "   Фантастичні істоти (химери, грифони), моралізаторська програма",
            "",
            "Д) З НАТУРАЛІСТИЧНИМИ РОСЛИНАМИ → фото 1 — ГОТИЧНИЙ",
            "   Дубове листя, плющ — спостереження з натури",
            "   Готика повернулась до живої природи (XIII–XIV ст.)",
        ],
        [],
        "Романські капітелі = умовність; Готичні = природність і деталізація"
    )

    # ── Q12 ────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 12,
        "Монументальне мистецтво середньовічних храмів",
        [
            "А) «Оранта» → фото 3",
            "   Київська Русь · Собор Святої Софії, Київ · МОЗАЇКА",
            "   Богоматір з піднятими руками — заступниця міста",
            "",
            "Б) «Княжі доньки» → фото 5",
            "   Київська Русь · Собор Святої Софії · ФРЕСКА",
            "   Процесія доньок Ярослава Мудрого",
            "",
            "В) «Ангел, що звиває небо» → фото 6",
            "   Київська Русь · Церква Спаса на Нередиці · ФРЕСКА",
            "",
            "Г) «Богородиця з Немовлям» → фото 2",
            "   Франція · Шартрський собор · ВІТРАЖ",
            "",
            "Д) «Діви мудрі та нерозумні» → фото 4",
            "   Франція чи Німеччина · портал собору · СТАТУЇ",
            "",
            "Е) «Благовіщення» → фото 1",
            "   Франція · Реймський собор · СТАТУЇ (кам'яний рельєф порталу)",
        ],
        ["byzantine_mosaic"],
        "Техніки: мозаїка = Візантія/Русь; вітраж і статуї = Захід"
    )

    # ── Q13 ────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 13,
        "Характерні ознаки візантійських зображень",
        [
            "✅  ПРАВИЛЬНІ ВІДПОВІДІ: Б, В, Д, Е, Ж, И",
            "",
            "Б) ВЕЛИКІ ОЧІ — «вікна душі», духовна зосередженість",
            "В) СПЛОЩЕНЕ ТІЛО — відхід від тілесності, перевага духовного",
            "Д) УМОВНИЙ ЗОЛОТИЙ ФОН — «Боже світло», позачасовий простір",
            "Е) ОБЕРНЕНА ПЕРСПЕКТИВА — об'єкти більшають з глибиною",
            "   (навпаки від природного сприйняття — глядач у центрі)",
            "Ж) ФРОНТАЛЬНІ ЗОБРАЖЕННЯ — постаті обернені до глядача",
            "И) ЗАСТИГЛІСТЬ МІМІКИ — austeritas, духовна велич",
            "",
            "❌  НЕПРАВИЛЬНІ (ознаки РЕАЛІЗМУ, не Візантії):",
            "А) підкреслений об'єм тіла — це антична скульптура",
            "Г) перспективна глибина простору — це Ренесанс",
            "Є) профільні зображення — єгипетський канон",
            "З) виражена міміка — це елліністична скульптура",
        ],
        ["byzantine_mosaic"],
        "Мозаїка Равенни — еталон візантійського канону (V–VI ст.)"
    )

    # ── Q14 ────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 14,
        "Техніка з використанням смальти",
        [
            "✅  ПРАВИЛЬНІ ФОТО: 3, 4, 5",
            "✅  НАЗВА ТЕХНІКИ: МОЗАЇКА",
            "",
            "Смальта — непрозоре кольорове скло у вигляді кубиків/пластин",
            "Виготовлялась шляхом плавлення скла з металевими оксидами:",
            "  • золото → жовта/золотиста смальта",
            "  • мідь → зелена і синя",
            "  • кобальт → темно-синя",
            "",
            "Процес мозаїки:",
            "  1. Тонкий шар вапняного розчину (тиньк)",
            "  2. Вдавлення кубиків смальти під різними кутами",
            "  3. Кути дають ефект мерехтіння/мерцання при свічковому освітленні",
            "",
            "❌  НЕ СМАЛЬТА:",
            "фото 1 — перегородчаста емаль (cloisonné)",
            "фото 2 — виїмчаста емаль (champlevé)",
            "фото 6 — вітраж (кольорове скло у свинцевій оправі)",
        ],
        ["byzantine_mosaic"],
        "Мозаїка з смальтою — основна техніка оздоблення візантійських і давньоруських храмів"
    )

    # ── Q15 ────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 15,
        "Канонічні зображення Богородиці",
        [
            "А) «ОРАНТА» (Та, що молиться) → фото 2",
            "   Богородиця зображена фронтально, стоячи,",
            "   з піднятими та розведеними руками — жест молитви.",
            "   Без Немовляти. Приклад: Оранта Київської Софії",
            "",
            "Б) «ОДИГІТРІЯ» (Дороговказівниця) → фото 1",
            "   Тримає Немовля на лівій руці, права — вказує на Христа.",
            "   Пряма, строга постава, погляд вперед.",
            "   Приклад: Казанська ікона Богоматері",
            "",
            "В) «ЕЛЕУСА» (Ніжність / Замилування) → фото 6",
            "   Мати торкається щокою щоки Немовляти,",
            "   Ісус обіймає Богородицю рукою за шию.",
            "   Вираз взаємної ніжності і передчуття страстей.",
            "   Приклад: Володимирська ікона Богоматері (XII ст., Москва)",
            "",
            "Решта: фото 3 — Успіння; фото 4 — Мадонна (Леонардо); фото 5 — Благовіщення",
        ],
        ["orans_madonna", "hodegetria", "eleusa"],
        "Одигітрія = «Та, що вказує Шлях»; Елеуса = людяність і материнська любов"
    )

    # ── Q16 ────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 16,
        "Образи християнських сюжетів",
        [
            "А) ГРІХОПАДІННЯ → фото 2",
            "   Зображені: Адам і Єва, Змій-спокусник, Дерево пізнання добра і зла",
            "",
            "Б) РІЗДВО ХРИСТОВЕ → фото 3, 5",
            "   фото 3 — Різдво (романський рельєф): Марія, немовля Ісус, пастухи, Йосиф",
            "   фото 5 — Візантійська мозаїка Різдва: печера, ясла, зірка, волхви, ангели",
            "",
            "В) СВЯТА ЄВХАРИСТІЯ (Причастя) → фото 4",
            "   Зображені: Ісус Христос, апостоли, хліб і вино — встановлення причастя",
            "",
            "ІНШІ СЮЖЕТИ (не відповіді):",
            "фото 1 — Ангел (можливо Страшний суд або Апокаліпсис)",
            "фото 6 — Пантократор (Христос Вседержитель) у куполі",
            "фото 7 — Вхід Господній в Єрусалим (Вербна неділя)",
        ],
        [],
        "Євхаристія = центральний сакральний сюжет у апсиді візантійських храмів"
    )

    # ══════════════════════════════════════════════════════════════════════════
    section_slide(prs, "МИСТЕЦТВО ВІДРОДЖЕННЯ")

    # ── Q17 ────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 17,
        "Поширений біблійний герой у Флоренції XV–поч. XVI ст.",
        [
            "✅  ПРАВИЛЬНА ВІДПОВІДЬ: Б) ДАВИД",
            "",
            "Давид — символ свободи Флорентійської республіки,",
            "перемоги малого над великим (Голіаф = тирани).",
            "",
            "ЗОБРАЖЕННЯ ДАВИДА (всі три — флорентійські):",
            "",
            "фото 1 — Донателло, «Давид» (мармур), ≈1408–1409 рр.",
            "   Перша монументальна статуя Відродження",
            "",
            "фото 2 — Донателло, «Давид» (бронза), ≈1440–1443 рр.",
            "   Перша вільностояча оголена бронзова статуя після античності",
            "",
            "фото 3 — Андреа Верроккйо, «Давид» (бронза), ≈1473–1475 рр.",
            "",
            "фото 4 — Мікеланджело, «Давид» (мармур), 1501–1504 рр.",
            "   Шедевр: 5,17 м, Галерея Академії, Флоренція",
        ],
        ["david_michelangelo", "david_donatello", "david_verrocchio"],
        "Мікеланджелів «Давид» — вершина Ренесансу: ідеальне тіло + психологічна напруга"
    )

    # ── Q18 ────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 18,
        "Живописні техніки епохи Відродження",
        [
            "✅  ПРАВИЛЬНІ: А) темпера,  Б) олійний живопис,  Г) фреска",
            "",
            "А) ТЕМПЕРА → фото 2",
            "   Боттічеллі «Симонетта Веспуччі» (~1480)",
            "   Пігменти на яєчному жовтку або казеїні",
            "   Яскрава, але крихка; шар за шаром; швидко сохне",
            "",
            "Б) ОЛІЙНИЙ ЖИВОПИС → фото 3, 4",
            "   фото 3 — Леонардо «Мона Ліза» (~1503–1519)",
            "   фото 4 — Гірландайо «Портрет молодого чоловіка» (~1480)",
            "   Технологія нідерландців (ван Ейк) → Венеція → вся Європа",
            "   Глибина кольору, сфумато, можливість доопрацьовувати",
            "",
            "Г) ФРЕСКА → фото 1",
            "   Мікеланджело, Сікстинська капела (1508–1512)",
            "   Пігменти на вологій штукатурці — «buon fresco»",
            "   Нерозривне зрощення з поверхнею стіни",
            "",
            "❌  НЕ ВІДРОДЖЕННЯ: мозаїка (В), вітраж (Д), енкаустика (Е)",
        ],
        ["sistine_fresco", "botticelli_simonetta", "mona_lisa"],
        "Темпера → XIII–XV ст.; Олія → з XV ст.; Фреска → паралельно, особливо монументальна"
    )

    # ── Q19 ────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 19,
        "Твори з міфологічними персонажами",
        [
            "А) Венера, Марс, сатири → фото 6",
            "   Боттічеллі «Венера і Марс», ~1483",
            "",
            "Б) Венера, Зефір, Аура, ора → фото 2",
            "   Боттічеллі «Народження Венери», 1485",
            "   Венера виходить з моря на мушлі; Зефір і Аура дмуть; Ора вдягає її",
            "",
            "В) Венера, Амур, грації, Зефір, Флора, Меркурій → фото 3",
            "   Боттічеллі «Весна» (Primavera), 1482",
            "   9 персонажів: алегорія весни і неоплатонічної любові",
            "",
            "Г) Даная і Зевс → фото 4",
            "   Тіціан «Даная», ~1553 — Зевс зійшов у вигляді золотого дощу",
            "",
            "Д) Європа і Зевс → фото 5",
            "   Тіціан «Викрадення Європи», ~1562",
            "   Зевс перетворився на бика, Європа тримається за роги",
            "",
            "Е) Венера і Амур → фото 1",
            "   Тіціан «Венера з дзеркалом», ~1555",
        ],
        ["botticelli_venus", "botticelli_spring", "titian_europa", "botticelli_venus_mars"],
        "Боттічеллі повернув античні сюжети у велике мистецтво вперше після 1000-річної перерви"
    )

    # ── Q20 ────────────────────────────────────────────────────────────────────
    answer_slide_1col(prs, 20,
        "Назви творів та імена художників Відродження",
        [
            "А) Андреа Мантенья → фото 4",
            "   Картина: «Мертвий Христос» (Lamento sul Cristo morto), ~1480–1490",
            "   Радикальне скорочення у перспективі (ноги на глядача)",
            "",
            "Б) «Портрет герцога Федеріго…» → фото 2",
            "   Автор: П'єро делла Франческа, 1473–1475",
            "   Диптих: Урбіно, профільні портрети на тлі пейзажу",
            "",
            "В) Джованні Белліні → фото 6",
            "   Картина: «Мадонна на лузі» (Madonna del Prato), ~1505",
            "   Лірична Мадонна з сонним Немовлям, відкритий пейзаж",
            "",
            "Г) «Венера Урбінська» → фото 5",
            "   Автор: Тіціан, 1538 — перша лежача Венера в масляному живописі",
            "",
            "Д) Фра Анджеліко → фото 1",
            "   Вівтарна картина: «Благовіщення», 1432",
            "   Ліворуч — вигнання з Раю (Адам і Єва), праворуч — Архангел і Марія",
            "",
            "Е) Мазаччо → фото 3",
            "   Фреска: «Данина» (Tribute Money), ~1425, капела Бранкаччі, Флоренція",
        ],
        ["fra_angelico_annunciation", "piero_diptych", "mantegna_dead_christ", "titian_venus_urbino"],
        "Мазаччо (†1428, 27 р.) — «батько живопису» Відродження: перша лінійна перспектива у фресці"
    )

    # ── FINAL SLIDE ────────────────────────────────────────────────────────────
    slide = blank_slide(prs)
    fill_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.08), GOLD)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), GOLD)
    add_text(slide, "Дякую за увагу", Inches(1), Inches(3.0), Inches(11.33), Inches(1.2),
             font_size=40, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, "Мистецтво давнього світу · Середні віки · Відродження",
             Inches(1), Inches(4.3), Inches(11.33), Inches(0.8),
             font_size=18, color=LGRAY, align=PP_ALIGN.CENTER)

    out_path = os.path.expanduser("~/Desktop/Art_History_Answers.pptx")
    prs.save(out_path)
    print(f"\n✅  Presentation saved: {out_path}")
    print(f"   Slides: {len(prs.slides)}")

if __name__ == "__main__":
    build()
